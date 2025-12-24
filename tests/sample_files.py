from __future__ import annotations

import csv
import io
import json
import math
import struct
import wave
from pathlib import Path


def create_sample_files(base_dir: Path) -> dict[str, Path]:
    base_dir.mkdir(parents=True, exist_ok=True)
    files: dict[str, Path] = {}

    html_path = base_dir / "sample.html"
    html_path.write_text(
        "<!doctype html><html><body><h1>Hello HTML</h1></body></html>",
        encoding="utf-8",
    )
    files[".html"] = html_path

    csv_path = base_dir / "sample.csv"
    with csv_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["title", "value"])
        writer.writerow(["Hello CSV", "42"])
    files[".csv"] = csv_path

    json_path = base_dir / "sample.json"
    json_path.write_text(json.dumps({"message": "Hello JSON"}), encoding="utf-8")
    files[".json"] = json_path

    xml_path = base_dir / "sample.xml"
    xml_path.write_text("<root><message>Hello XML</message></root>", encoding="utf-8")
    files[".xml"] = xml_path

    _write_simple_pdf(base_dir / "sample.pdf", "Hello PDF")
    files[".pdf"] = base_dir / "sample.pdf"

    _write_docx(base_dir / "sample.docx")
    files[".docx"] = base_dir / "sample.docx"

    _write_pptx(base_dir / "sample.pptx")
    files[".pptx"] = base_dir / "sample.pptx"

    _write_xlsx(base_dir / "sample.xlsx")
    files[".xlsx"] = base_dir / "sample.xlsx"

    _write_images(base_dir, files)

    wav_path = base_dir / "sample.wav"
    _write_wav(wav_path)
    files[".wav"] = wav_path

    mp3_path = base_dir / "sample.mp3"
    if _try_write_mp3(mp3_path, wav_path):
        files[".mp3"] = mp3_path

    return files


def _write_simple_pdf(path: Path, text: str) -> None:
    content = f"BT\n/F1 24 Tf\n72 120 Td\n({text}) Tj\nET\n"
    objects = [
        "1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        "2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        "3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 200] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n",
        f"4 0 obj\n<< /Length {len(content)} >>\nstream\n{content}endstream\nendobj\n",
        "5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
    ]

    buffer = io.BytesIO()
    buffer.write(b"%PDF-1.4\n")
    offsets = [0]

    for obj in objects:
        offsets.append(buffer.tell())
        buffer.write(obj.encode("latin-1"))

    xref_offset = buffer.tell()
    buffer.write(f"xref\n0 {len(offsets)}\n".encode("latin-1"))
    buffer.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        buffer.write(f"{offset:010d} 00000 n \n".encode("latin-1"))

    buffer.write(
        (
            "trailer\n"
            f"<< /Size {len(offsets)} /Root 1 0 R >>\n"
            "startxref\n"
            f"{xref_offset}\n"
            "%%EOF\n"
        ).encode("latin-1")
    )

    path.write_bytes(buffer.getvalue())


def _write_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_heading("Hello DOCX", level=1)
    document.add_paragraph("Hello from DOCX")
    document.save(path)


def _write_pptx(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello PPTX"
    subtitle.text = "Hello from PPTX"
    presentation.save(path)


def _write_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    worksheet = workbook.active
    worksheet["A1"] = "Hello XLSX"
    worksheet["B1"] = "42"
    workbook.save(path)


def _write_images(base_dir: Path, files: dict[str, Path]) -> None:
    from PIL import Image, ImageDraw

    image = Image.new("RGB", (120, 60), color=(255, 255, 255))
    draw = ImageDraw.Draw(image)
    draw.text((10, 20), "Hi", fill=(0, 0, 0))

    png_path = base_dir / "sample.png"
    image.save(png_path, format="PNG")
    files[".png"] = png_path

    jpg_path = base_dir / "sample.jpg"
    image.save(jpg_path, format="JPEG")
    files[".jpg"] = jpg_path

    jpeg_path = base_dir / "sample.jpeg"
    image.save(jpeg_path, format="JPEG")
    files[".jpeg"] = jpeg_path

    gif_path = base_dir / "sample.gif"
    image.save(gif_path, format="GIF")
    files[".gif"] = gif_path


def _write_wav(path: Path, *, duration: float = 1.0, frequency: float = 440.0) -> None:
    framerate = 44100
    total_frames = int(duration * framerate)

    with wave.open(str(path), "w") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(framerate)

        frames = bytearray()
        for i in range(total_frames):
            sample = int(32767 * math.sin(2 * math.pi * frequency * i / framerate))
            frames.extend(struct.pack("<h", sample))

        wav_file.writeframes(frames)


def _try_write_mp3(path: Path, wav_path: Path) -> bool:
    try:
        from pydub import AudioSegment
        from pydub.utils import which
    except Exception:
        return False

    if which("ffmpeg") is None:
        return False

    try:
        audio = AudioSegment.from_wav(wav_path)
        audio.export(path, format="mp3")
    except Exception:
        return False

    return True
